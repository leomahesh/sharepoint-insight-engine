import os
import pandas as pd
from pydantic import BaseModel
from typing import Optional, List
import logging

# Simple parsers
# For production, 'unstructured' is better but requires more deps.
# We will use pypdf for PDF and pandas for Excel.
# For Docx we might need python-docx, but let's see if we can do basic text read or add dep.

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

logger = logging.getLogger(__name__)

def extract_text_from_file(file_path: str, file_type: str) -> str:
    """
    Extracts text content from a file based on its type.
    """
    content = ""
    try:
        if file_type == 'pdf':
            if not PdfReader:
                return "[Error] pypdf not installed."
            reader = PdfReader(file_path)
            for page in reader.pages:
                content += page.extract_text() + "\n"
        
        elif file_type in ['xlsx', 'xls', 'csv']:
            if file_type == 'csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            # Convert full dataframe to string representation
            content = df.to_string()
        
        elif file_type in ['docx', 'doc']:
            if not DocxDocument:
                 return "[Error] python-docx not installed."
            doc = DocxDocument(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])

        elif file_type in ['pptx', 'ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                content = "\n".join(text_runs)
            except ImportError:
                 return "[Error] python-pptx not installed."
        
        elif file_type in ['pptx', 'ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                content = "\n".join(text_runs)
            except ImportError:
                 return "[Error] python-pptx not installed."
        
        elif file_type in ['txt', 'md', 'json']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        
        else:
            # Fallback for unknown
            content = f"[Preview] File type {file_type} text extraction not fully supported yet."

    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        content = f"[Error reading file: {str(e)}]"

    return content
